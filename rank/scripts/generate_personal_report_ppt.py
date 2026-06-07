# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE = r"c:\Users\12113\Desktop\mid-term report.pptx"
OUTPUT = r"c:\Users\12113\Desktop\OJ排行榜模块-个人报告.pptx"


def set_text(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line


def at(slide, idx, text):
    if idx < len(slide.shapes) and slide.shapes[idx].has_text_frame:
        set_text(slide.shapes[idx], text)


def add_content_box(slide, title, body):
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    if p0.runs:
        p0.runs[0].font.bold = True
        p0.runs[0].font.size = Pt(22)
    for line in body.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        if p.runs:
            p.runs[0].font.size = Pt(15)


def main():
    prs = Presentation(TEMPLATE)

    # Slide 0 封面
    s = prs.slides[0]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        if "期中" in t or ("网站" in t and len(t) < 20):
            set_text(shape, "在线评测系统 · 个人报告")
        elif "OJ website" in t or "mid-term" in t.lower():
            set_text(shape, "OJ Online Judge System — Personal Report")
        elif "XXX" in t or "小组成员" in t:
            set_text(shape, "负责模块：排行榜模块（rank）\n姓名：[请填写]    学号：[请填写]")

    # Slide 1 引言
    for shape in prs.slides[1].shapes:
        if shape.has_text_frame and len(shape.text_frame.text) > 40:
            set_text(
                shape,
                "本人负责 OJ 系统排行榜模块（OJ/rank）的设计、开发与交付。\n"
                "模块提供全球榜、比赛榜、用户题目统计及 LeetCode 风格可视化前端，\n"
                "并完成 MySQL 持久化、Windows 运行包与模块文档，支撑团队联调。",
            )

    # Slide 2 目录
    for shape in prs.slides[2].shapes:
        if shape.has_text_frame and "CONTENT" in shape.text_frame.text:
            set_text(shape, "01 个人任务\n02 理论与技术\n03 工具与辅助\n04 工作量与总结")

    # Slide 3 章节01
    s = prs.slides[3]
    at(s, 6, "个人任务")
    at(s, 7, "01")

    # Slide 4 个人任务
    s = prs.slides[4]
    at(s, 4, "个人任务")
    at(s, 9, "全球榜 / 比赛榜 / 用户统计\nREST API / CSV 导出 / 健康检查")
    at(s, 10, "核心功能开发")
    at(s, 11, "工程化交付")
    at(s, 14, "主要交付物")
    at(s, 15, "排行榜模块从 0 到可联调交付：\n业务逻辑 + HTTP 服务 + 前端 + 持久化")
    at(s, 16, "LeetCode 风格前端（ECharts 柱状图）\nMySQL 持久化 + Windows 运行包\n《排行榜模块使用文档》约 370 行")

    # Slide 5 模块职责
    s = prs.slides[5]
    at(s, 1, "排行榜模块职责")
    at(s, 4, "分页查询、多关键字排序\nGET /api/leaderboard/global")
    at(s, 5, "全球排行榜")
    at(s, 7, "按比赛 ID 查询排名\nGET /api/leaderboard/contest/{id}")
    at(s, 8, "比赛排行榜")
    at(s, 10, "Easy/Medium/Hard 完成统计\nGET /api/leaderboard/users/{id}/stats")
    at(s, 11, "用户题目统计")
    at(s, 24, "个人任务")
    at(s, 27, "数据持久化")

    # Slide 6
    s = prs.slides[6]
    at(s, 3, "个人任务")

    # Slide 7 章节02
    s = prs.slides[7]
    at(s, 6, "理论与技术")
    at(s, 7, "02")

    # Slide 8 软件工程理论
    s = prs.slides[8]
    at(s, 3, "理论与技术")
    at(s, 7, "分层架构")
    at(s, 10, "设计模式")
    at(s, 12, "架构分层")
    at(s, 13, "Handler → API → Service → Repository\n仓储模式 + 工厂模式（内存/MySQL 切换）\nRESTful API、事务处理、配置外部化")

    # Slide 9 模块分层
    s = prs.slides[9]
    at(s, 3, "理论与技术")
    at(s, 14, "前端 web/")
    at(s, 15, "排行榜 API")
    at(s, 16, "Service 层")
    at(s, 17, "Storage 层")
    at(s, 22, "HTML + Tailwind + ECharts\nLeetCode 风格领奖台与柱状图")
    at(s, 23, "全球榜/比赛榜/统计/CSV 导出\n/health 健康检查接口")
    at(s, 24, "LeaderboardService 业务逻辑\nStatsAnalyzer 百分位洞察")
    at(s, 25, "InMemory / MysqlLeaderboardRepository\nlibmariadbclient C API 持久化")

    # Slides 10-12 补充页
    add_content_box(
        prs.slides[10],
        "01 软件工程理论应用",
        "- 分层架构：Handler→API→Service→Repository 职责分离\n"
        "- 仓储模式：ILeaderboardRepository 抽象存储实现\n"
        "- 工厂模式：create_leaderboard_repository() 按环境自动选型\n"
        "- 依赖倒置：Service 依赖抽象接口，不绑定具体存储",
    )
    add_content_box(
        prs.slides[11],
        "02 人工智能相关",
        "排行榜为传统业务模块，AI 能力由讨论区（Gemini）承担。\n"
        "本人预留数据接口供后续 AI 应用：\n"
        "- 用户百分位洞察 GET /users/{id}/insight\n"
        "- 全局统计摘要 GET /global/summary\n"
        "- 可对接学习路径推荐、智能排名分析",
    )
    add_content_box(
        prs.slides[12],
        "03 课程知识总结",
        "- RESTful API 设计与 JSON 响应规范\n"
        "- MySQL 表设计、外键约束（Users 表）、事务一致性\n"
        "- 单元测试 + MySQL 集成测试（连接失败优雅跳过）\n"
        "- 配置外部化：OJ_MYSQL_* 环境变量管理",
    )

    # Slide 13 章节03
    s = prs.slides[13]
    at(s, 6, "工具与辅助")
    at(s, 7, "03")

    # Slide 14 技术与工具
    s = prs.slides[14]
    at(s, 4, "工具与辅助")
    at(s, 8, "C++17 + CMake + Ninja\ncommon HTTP/日志 + ECharts + Tailwind")
    at(s, 9, "学习了哪些技术？")
    at(s, 11, "使用了哪些辅助工具？")
    at(s, 12, "Cursor IDE（AI 辅助编程）")
    at(s, 13, "Git + GitHub + PowerShell")
    at(s, 14, "视频演示 Cursor + Git（≤5 分钟）")

    # Slide 15 技术栈详情
    s = prs.slides[15]
    at(s, 3, "技术栈详情")
    at(s, 14, "构建工具")
    at(s, 15, "版本管理")
    at(s, 16, "CMake + Ninja + MSYS2\npacman install libmariadbclient")
    at(s, 17, "后端技术")
    at(s, 18, "C++17、oj::HttpServer、oj::Router\n分层架构 + 单元测试")
    at(s, 19, "数据存储")
    at(s, 20, "团队 myOJ 库\nlibmariadbclient C API 持久化")

    # Slide 16 辅助工具实践
    s = prs.slides[16]
    at(s, 4, "辅助工具实践")
    at(s, 15, "Cursor IDE\n参考 discussion 实现 MySQL 持久化\nCrow → common HTTP 迁移")
    at(s, 16, "需求分析")
    at(s, 18, "Git + GitHub\nfeat/rank-runtime-package 分支\nfork 推送 + PR 合并")
    at(s, 19, "版本协作")
    at(s, 21, "PowerShell")
    at(s, 22, "sync_runtime.ps1 自动化\n运行包同步与一键启动")
    at(s, 24, "浏览器 DevTools")
    at(s, 25, "调试 REST API\n验证 /health 与排行榜数据\n（视频演示 ≤5 分钟）")

    # Slide 17 章节04
    s = prs.slides[17]
    at(s, 6, "工作量与总结")
    at(s, 7, "04")

    # Slide 18 工作量
    s = prs.slides[18]
    at(s, 4, "工作量与总结")
    at(s, 7, "个人工作量统计")
    at(s, 12, "代码约 3000+ 行\nC++ 1900+ / JS 578 / HTML+CSS 285\n文档约 370 行 / 缺陷修复 7 项")
    at(s, 13, "项目管理参与度\nGit 分支 + fork/PR 协作\nCMake 配置 + runtime_package 同步\n单元测试 + 集成测试 + /health")

    # Slide 19 完成与待完善
    s = prs.slides[19]
    at(s, 3, "工作量与总结")
    at(s, 11, "已完成工作")
    at(s, 12, "排行榜业务 + REST API\nLeetCode 前端 + ECharts 可视化\nMySQL 持久化（C API）+ 运行包\n模块文档 + 扩展测试")
    at(s, 14, "待完善工作")
    at(s, 15, "judge 判题结果自动写入排行榜\n纳入根工程 CMakeLists 统一构建\n与团队 MySQL 服务器联调验证")

    # Slide 20 后续与时间
    s = prs.slides[20]
    at(s, 3, "工作量与总结")
    at(s, 15, "数据对接")
    at(s, 16, "统一构建")
    at(s, 17, "事件驱动")
    at(s, 18, "汇报时间")
    at(s, 22, "时间控制建议")
    at(s, 28, "judge → rank on_submission()\n判题结果自动更新排行榜")
    at(s, 29, "rank 加入 OJ/CMakeLists.txt\n团队一键编译全部模块")
    at(s, 30, "完善 /health 运维监控\n生产环境部署验证")
    at(s, 31, "建议 8-10 分钟汇报：\n任务 1.5 + 理论 2 + 工具 2 + 工作量 2 + 总结 1 min")

    # Slide 21 致谢
    s = prs.slides[21]
    at(s, 4, "谢谢观看")
    at(s, 6, "日期：2026/6/6")

    prs.save(OUTPUT)
    print("Saved:", OUTPUT)


if __name__ == "__main__":
    main()
