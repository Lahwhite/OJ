# -*- coding: utf-8 -*-
from pptx import Presentation

path = r"c:\Users\12113\Desktop\OJ排行榜模块-个人报告.pptx"
prs = Presentation(path)
lines = []
for si in range(len(prs.slides)):
    texts = []
    for shape in prs.slides[si].shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            texts.append(shape.text_frame.text.strip())
    if texts:
        lines.append(f"=== Slide {si} ===")
        lines.extend(texts)
        lines.append("")

with open(r"f:\软工实验（）新\OJ\rank\scripts\ppt_content.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(lines))
print("written")
